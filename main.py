import time
import os
import sys
import winreg
import getpass
import hashlib
import string
from datetime import datetime
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import threading
import ctypes
import xml.etree.ElementTree as ET
import urllib.request
import urllib.parse

# Load environment variables from .env file (supports python-dotenv or manual parsing fallback)
try:
    from dotenv import load_dotenv
    load_dotenv(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env"))
except ImportError:
    _env_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env")
    if os.path.exists(_env_path):
        with open(_env_path, "r", encoding="utf-8") as f:
            for line in f:
                line_stripped = line.strip()
                if line_stripped and not line_stripped.startswith("#") and "=" in line_stripped:
                    try:
                        key, val = line_stripped.split("=", 1)
                        os.environ[key.strip()] = val.strip().strip('"').strip("'")
                    except Exception:
                        pass

def is_admin():
    try:
        return ctypes.windll.shell32.IsUserAnAdmin() == 1
    except:
        return False

# ── Optional: pynput + win32gui (real-time MODIFIED & OPENED) ─────────────
try:
    from pynput import keyboard as pynput_keyboard
    import win32gui
    import win32event
    import win32api
    import winerror
    HAS_REALTIME = True
except ImportError:
    HAS_REALTIME = False
    print("[!] pynput/pywin32 not installed  →  pip install pynput pywin32")

# ── Optional: psutil (SMB network monitor) ─────────────
try:
    import psutil
    HAS_PSUTIL = True
except ImportError:
    HAS_PSUTIL = False
    print("[!] psutil not installed  →  pip install psutil")


# ================== PATH SETUP ==================
BASE_DIR       = os.path.dirname(os.path.abspath(__file__))
HONEY_DIR      = os.path.join(BASE_DIR, "honeyfiles")
HONEY_FILES = {
    "passwords.txt": os.path.join(HONEY_DIR, "passwords.txt"),
    "confidential.docx": os.path.join(HONEY_DIR, "confidential.docx"),
    "financial_budget_2024.xlsx": r"D:\financial_budget_2024.xlsx",
    "employee_salaries.pdf": r"D:\employee_salaries.pdf",
    "private_keys.pem": r"D:\private_keys.pem",
    "client_database.csv": r"D:\client_database.csv",
    "merger_plans.pptx": r"D:\merger_plans.pptx",
    "source_code_backup.zip": r"D:\source_code_backup.zip",
    "admin_credentials.xml": r"D:\admin_credentials.xml",
    "payroll_2025.xlsx": r"D:\payroll_2025.xlsx"
}
LOG_FILE       = os.path.join(BASE_DIR, "logs", "alerts.log")

# ================== RISK CONFIG ==================
RISK_SCORES = {
    "OPENED":   10,
    "MODIFIED": 40,
    "MOVED":    60,
    "COPIED":   70,
    "DELETED":  90,
    "RENAMED":  90,
    "SMB_GET":  80,
    "SMB_PUSH": 90,
    "SMB_ACCESS": 50,
}

# ================== ALERT CONFIG ==================
# Telegram settings (loaded dynamically from environment variables for security)
TELEGRAM_BOT_TOKEN = os.environ.get("TELEGRAM_BOT_TOKEN", "your_bot_token")
TELEGRAM_CHAT_ID = os.environ.get("TELEGRAM_CHAT_ID", "your_chat_id")

ENABLE_ALERTS = True  # Alerts are now enabled



# ================== GLOBALS ==================
copied_candidates = {}   # {abspath: timestamp}
known_copies      = set()
ORIGINAL_HASHES   = {}
last_hashes       = {}
active_smb_ips    = set()
file_open_states  = {}   # tracking open durations: { file_name: time_opened_timestamp }

# Debounce flags so we don't spam logs
_last_opened_log  = 0.0
_last_typing_log  = 0.0
_OPEN_DEBOUNCE    = 3.0
_TYPE_DEBOUNCE    = 2.0


# ================== UTIL FUNCTIONS ==================
def send_critical_alert(severity, action, file_name, location, smb_info):
    if not ENABLE_ALERTS:
        return
        
    subject = f"CRITICAL THREAT: Honeyfile '{file_name}' Accessed!"
    body = f"""
SEVERITY: {severity}
ACTION: {action}
TARGET FILE: {file_name}
LOCATION: {location}

DETAILS:
{smb_info}
Please check the system immediately.
"""
    import html
    
    # Escape variables to prevent HTML parsing errors on Telegram's side
    safe_file_name = html.escape(str(file_name))
    safe_action = html.escape(str(action))
    safe_severity = html.escape(str(severity))
    safe_location = html.escape(str(location))
    safe_smb_info = html.escape(str(smb_info.strip())) if smb_info.strip() else "Local (No remote SMB interaction detected)."
    
    # Send Telegram Message
    telegram_html = f"""
🚨 <b>CRITICAL INTRUSION DETECTED</b> 🚨

Someone has compromised a restricted honeyfile on your system!

📝 <b>Event Details:</b>
• <b>Target File:</b> <code>{safe_file_name}</code>
• <b>Action Taken:</b> <tg-spoiler>{safe_action}</tg-spoiler>
• <b>Severity Level:</b> <b>{safe_severity}</b>

📂 <b>Absolute Path:</b>
<code>{safe_location}</code>

🕵️‍♂️ <b>Network & Actor Details:</b>
<blockquote>{safe_smb_info}</blockquote>

<i>Check your system immediately!</i>
"""
    try:
        if TELEGRAM_BOT_TOKEN != "your_bot_token" and TELEGRAM_CHAT_ID != "your_chat_id":
            telegram_url = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}/sendMessage"
            data = urllib.parse.urlencode({
                'chat_id': TELEGRAM_CHAT_ID,
                'text': telegram_html,
                'parse_mode': 'HTML'
            }).encode('utf-8')
            
            req = urllib.request.Request(telegram_url, data=data)
            with urllib.request.urlopen(req, timeout=10) as response:
                if response.getcode() == 200:
                    print(f"\n[+] Critical alert securely sent to your Telegram.")
        else:
            print("\n[-] Telegram alert skipped (Token/Chat ID not configured yet).")
    except Exception as e:
        print(f"\n[!] Failed to send Telegram alert: {e}")
        # Log the error to a file so we can see it even if running in pythonw.exe
        try:
            error_log_path = os.path.join(BASE_DIR, "logs", "error.log")
            with open(error_log_path, "a") as ef:
                ef.write(f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] Telegram Error: {str(e)}\n")
        except:
            pass
            
        # Fallback: Try sending without HTML parse_mode in case it was a formatting rejection
        try:
            fallback_data = urllib.parse.urlencode({
                'chat_id': TELEGRAM_CHAT_ID,
                'text': f"CRITICAL INTRUSION DETECTED\nFile: {file_name}\nAction: {action}\nLocation: {location}\nSeverity: {severity}\nInfo: {smb_info}"
            }).encode('utf-8')
            req = urllib.request.Request(telegram_url, data=fallback_data)
            urllib.request.urlopen(req, timeout=10)
        except:
            pass

def get_file_hash(path):
    try:
        with open(path, "rb") as f:
            return hashlib.md5(f.read()).hexdigest()
    except Exception:
        return None

def calculate_severity(score):
    if score < 20: return "LOW"
    elif score < 40: return "HIGH"
    else: return "CRITICAL"

def log_event(action, location, force_smb_info=None, file_name=None):
    global _last_opened_log, _last_typing_log, file_open_states

    username  = getpass.getuser()
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    fname = file_name if file_name else os.path.basename(location)

    # Track logic for OPEN duration
    if action in ["OPENED", "SMB_GET"]:
        if fname not in file_open_states:
            file_open_states[fname] = time.time()
            
    if action == "CLOSED":
        if fname in file_open_states:
            duration = time.time() - file_open_states[fname]
            del file_open_states[fname]
        else:
            duration = 0 # Missed OPEN event
            
        if duration <= 60:
            risk_score = 10     # Low (< 1 min)
        elif duration <= 300:
            risk_score = 30     # High / Medium (1 to 5 min)
        else:
            risk_score = 60     # Critical (> 5 min)
    else:
        risk_score = RISK_SCORES.get(action, 50)

    severity   = calculate_severity(risk_score)

    smb_info = ""
    if force_smb_info:
        smb_info = f"Remote Details: {force_smb_info}\n"

    log_message = f"""
[{timestamp}]
User     : {username}
Action   : {action}
File     : {fname}
Location : {location}
{smb_info}Severity : {severity}
Risk Score: {risk_score}
-----------------------------------------
"""
    print(log_message)

    try:
        with open(LOG_FILE, "a") as f:
            f.write(log_message)
            f.flush()
            os.fsync(f.fileno())
    except Exception as e:
        print(f"[!] Failed to write to {LOG_FILE}: {e}")

    # Trigger Telegram alert on critical severity
    if severity == "CRITICAL":
        threading.Thread(target=send_critical_alert, args=(severity, action, fname, location, smb_info), daemon=True).start()


# ================== EVENT HANDLER (watchdog) ==================
class HoneyFileHandler(FileSystemEventHandler):

    def get_tracked_file(self, path):
        if not path: return None
        lower_path = path.lower()
        if lower_path.endswith('.lnk'): return None
        for f in HONEY_FILES.keys():
            if f.lower() in lower_path:
                return f
        return None

    def on_modified(self, event):
        global last_hashes
        if event.is_directory: return
        
        tracked_file = self.get_tracked_file(event.src_path)
        if not tracked_file: return
        
        abs_path = os.path.abspath(event.src_path)
        norm_abs_path = os.path.normcase(abs_path)
        original_path = os.path.abspath(HONEY_FILES[tracked_file])
        norm_orig_path = os.path.normcase(original_path)

        if norm_abs_path == norm_orig_path or norm_abs_path in known_copies:
            if norm_abs_path in copied_candidates:
                if time.time() - copied_candidates[norm_abs_path] < 5:
                    return

            current_hash = get_file_hash(abs_path)
            if current_hash is None: return

            if current_hash != last_hashes.get(tracked_file):
                last_hashes[tracked_file] = current_hash
                log_event("MODIFIED", event.src_path, file_name=tracked_file)
                if norm_abs_path in copied_candidates:
                    del copied_candidates[norm_abs_path]

        else:
            if norm_abs_path not in copied_candidates:
                copied_candidates[norm_abs_path] = time.time()

            current_hash = get_file_hash(abs_path)
            if ORIGINAL_HASHES.get(tracked_file) and current_hash == ORIGINAL_HASHES[tracked_file]:
                log_event("COPIED", event.src_path, file_name=tracked_file)
                known_copies.add(norm_abs_path)
                if norm_abs_path in copied_candidates:
                    del copied_candidates[norm_abs_path]

    def on_deleted(self, event):
        tracked_file = self.get_tracked_file(event.src_path)
        if tracked_file:
            abs_path = os.path.abspath(event.src_path)
            norm_abs_path = os.path.normcase(abs_path)
            original_path = os.path.abspath(HONEY_FILES[tracked_file])
            norm_orig_path = os.path.normcase(original_path)

            if norm_abs_path == norm_orig_path or norm_abs_path in known_copies:
                log_event("DELETED", event.src_path, file_name=tracked_file)
                known_copies.discard(norm_abs_path)
                copied_candidates.pop(norm_abs_path, None)
            else:
                # Its a temporary file or lock file getting deleted, meaning file was closed.
                log_event("CLOSED", event.src_path, file_name=tracked_file)

    def on_moved(self, event):
        tracked_src = self.get_tracked_file(event.src_path)
        if tracked_src:
            src_dir = os.path.dirname(os.path.abspath(event.src_path))
            dest_dir = os.path.dirname(os.path.abspath(event.dest_path))
            if src_dir == dest_dir:
                log_event("RENAMED", f"{event.src_path} -> {event.dest_path}", file_name=tracked_src)
            else:
                log_event("MOVED", f"{event.src_path} -> {event.dest_path}", file_name=tracked_src)

    def on_created(self, event):
        tracked_file = self.get_tracked_file(event.src_path)
        if tracked_file:
            abs_path = os.path.abspath(event.src_path)
            norm_abs_path = os.path.normcase(abs_path)
            original_path = os.path.abspath(HONEY_FILES[tracked_file])
            if norm_abs_path == os.path.normcase(original_path):
                return
            
            fname = os.path.basename(norm_abs_path)
            if fname.startswith("~$") or fname.startswith(".~") or fname.endswith(".tmp"):
                log_event("OPENED", event.src_path, file_name=tracked_file)
                
            copied_candidates[norm_abs_path] = time.time()

# ================== FILE OPEN DETECTOR (Window Tracker) ==================
def monitor_file_open():
    global _last_opened_log
    if not HAS_REALTIME: return

    print("[*] File-open detector active (win32gui Window Tracker).")
    # Mapping "passwords" -> "passwords.txt", "confidential" -> "confidential.docx"
    honey_keywords = {os.path.splitext(f)[0].lower(): f for f in HONEY_FILES.keys()}
    last_window_title = ""

    while True:
        try:
            hwnd  = win32gui.GetForegroundWindow()
            title = win32gui.GetWindowText(hwnd).lower()

            for kw, full_name in honey_keywords.items():
                if kw in title and title != last_window_title:
                    now = time.time()
                    if now - _last_opened_log > _OPEN_DEBOUNCE:
                        _last_opened_log = now
                        log_event("OPENED", HONEY_FILES[full_name], file_name=full_name)
                    break
            last_window_title = title
        except Exception:
            pass
        time.sleep(0.5)

# ================== REAL-TIME TYPING DETECTOR (pynput) ==================
def monitor_realtime_typing():
    global _last_typing_log
    if not HAS_REALTIME: return
    print("[*] Real-time typing detector active (pynput).")

    honey_keywords = {os.path.splitext(f)[0].lower(): f for f in HONEY_FILES.keys()}

    def on_key_press(key):
        global _last_typing_log
        try:
            hwnd  = win32gui.GetForegroundWindow()
            title = win32gui.GetWindowText(hwnd).lower()
            for kw, full_name in honey_keywords.items():
                if kw in title:
                    now = time.time()
                    if now - _last_typing_log > _TYPE_DEBOUNCE:
                        _last_typing_log = now
                        log_event("MODIFIED", HONEY_FILES[full_name], file_name=full_name)
                    break
        except Exception:
            pass

    with pynput_keyboard.Listener(on_press=on_key_press) as listener:
        listener.join()

# ================== COPY PROCESSOR ==================
def process_copy_events():
    while True:
        time.sleep(2)
        current_time = time.time()
        for path in list(copied_candidates.keys()):
            t = copied_candidates.get(path)
            if t is None: continue

            if current_time - t > 2:
                copied_hash = get_file_hash(path)
                
                # Check against ALL known hashes to see if it's a copy of any honeyfile
                matched_file = None
                for fname, ohash in ORIGINAL_HASHES.items():
                    if ohash and copied_hash == ohash:
                        matched_file = fname
                        break

                if matched_file:
                    original_path = os.path.abspath(HONEY_FILES[matched_file])
                    # Ensure we don't log the original file due to some race condition/case mismatch
                    if path == os.path.normcase(original_path):
                        copied_candidates.pop(path, None)
                        continue

                    if path not in known_copies:
                        log_event("COPIED", path, file_name=matched_file)
                        known_copies.add(path)
                    copied_candidates.pop(path, None)
                elif current_time - t > 10:
                    copied_candidates.pop(path, None)

# ================== DRIVE MONITOR ==================
def get_all_drives():
    drives = []
    for letter in string.ascii_uppercase:
        drive = f"{letter}:\\"
        if os.path.exists(drive):
            drives.append(drive)
    return drives

def monitor_drives(observer, handler):
    monitored = set()
    while True:
        for drive in set(get_all_drives()):
            if drive not in monitored:
                try:
                    observer.schedule(handler, path=drive, recursive=True)
                    monitored.add(drive)
                    print(f"[+] Now monitoring drive: {drive}")
                except Exception as e:
                    print(f"[!] Failed to monitor {drive}: {e}")
        time.sleep(5)

# ================== SMB / NETWORK MONITOR ==================
def monitor_smb_sessions():
    global active_smb_ips
    if not HAS_PSUTIL: return

    print("[*] SMB Network Monitor active. Monitoring port 445 for 'get' and 'push'.")
    while True:
        try:
            current_ips = set()
            for conn in psutil.net_connections(kind='tcp'):
                if conn.laddr and conn.laddr.port == 445 and conn.status == 'ESTABLISHED':
                    if conn.raddr:
                        ip = conn.raddr.ip
                        # Ignore local loopback connections which generate false positives
                        if ip not in ['127.0.0.1', '::1', '0.0.0.0', '::']:
                            current_ips.add(ip)
            new_ips = current_ips - active_smb_ips
            for ip in new_ips:
                print(f"[*] Network notice: Remote SMB Connection Established from {ip}")
            active_smb_ips = current_ips
        except Exception:
            pass
        time.sleep(2)

# ================== SECURITY LOG MONITOR (SMB_GET) ==================
def monitor_security_logs():
    if not is_admin(): return
    print("[*] Security Event Log Monitor active (Event ID 5145).")
    
    import subprocess
    # Enable Share Access Auditing System-wide
    subprocess.run(['auditpol', '/set', '/subcategory:Detailed File Share', '/success:enable', '/failure:enable'], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, creationflags=subprocess.CREATE_NO_WINDOW)
    
    last_time = datetime.utcnow().strftime('%Y-%m-%dT%H:%M:%S.000Z')
    processed_records = set()
    last_event_times = {} # debounce dictionary
    
    while True:
        try:
            query = f"*[System[(EventID=5145) and TimeCreated[@SystemTime>='{last_time}']]]"
            cmd = ['wevtutil', 'qe', 'Security', f'/q:{query}', '/e:Events']
            new_time = datetime.utcnow().strftime('%Y-%m-%dT%H:%M:%S.000Z')
            
            result = subprocess.run(cmd, capture_output=True, text=True, creationflags=subprocess.CREATE_NO_WINDOW)
            
            if result.stdout.strip():
                root = ET.fromstring(result.stdout)
                ns = {'ns': 'http://schemas.microsoft.com/win/2004/08/events/event'}
                
                for event in root.findall('ns:Event', ns):
                    system = event.find('ns:System', ns)
                    record_id = system.find('ns:EventRecordID', ns).text
                    
                    if record_id in processed_records: continue
                    processed_records.add(record_id)
                    if len(processed_records) > 2000: processed_records.pop()

                    event_data = event.find('ns:EventData', ns)
                    if event_data is not None:
                        obj_name = ""
                        ip_addr = ""
                        accesses = ""
                        for data in event_data.findall('ns:Data', ns):
                            name = data.get('Name')
                            if name == 'RelativeTargetName':
                                obj_name = data.text or ""
                            elif name == 'IpAddress':
                                ip_addr = data.text or ""
                            elif name == 'AccessList':
                                accesses = data.text or ""
                                
                        if obj_name:
                            matched_file = None
                            for fname in HONEY_FILES.keys():
                                if fname.lower() in obj_name.lower():
                                    matched_file = fname
                                    break
                            
                            # If they touched the dummy file over SMB, it's a hit.
                            if matched_file:
                                now = time.time()
                                if now - last_event_times.get(matched_file, 0) > 3.0:
                                    last_event_times[matched_file] = now
                                    log_event("SMB_GET", obj_name, force_smb_info=f"Read via Network from {ip_addr}", file_name=matched_file)

        except Exception as e:
            # print(e)
            pass
            
        last_time = new_time
        time.sleep(2)

# ================== STARTUP ==================
def add_to_startup():
    key = r"Software\Microsoft\Windows\CurrentVersion\Run"
    app_name = "HoneyFileMonitor"
    python_exe = sys.executable
    if python_exe.lower().endswith("python.exe"):
        pythonw = python_exe[:-10] + "pythonw.exe"
        if os.path.exists(pythonw): python_exe = pythonw
    script_path = os.path.abspath(__file__)
    command = f'"{python_exe}" "{script_path}"'
    try:
        registry_key = winreg.OpenKey(winreg.HKEY_CURRENT_USER, key, 0, winreg.KEY_SET_VALUE)
        winreg.SetValueEx(registry_key, app_name, 0, winreg.REG_SZ, command)
        winreg.CloseKey(registry_key)
        print("[+] Successfully registered to run at Windows Startup.")
    except Exception as e:
        print(f"[!] Failed to add to Startup: {e}")

def prevent_multiple_instances():
    if not HAS_REALTIME: return None
    mutex = win32event.CreateMutex(None, False, "Global\\HoneyFileMonitor_Mutex")
    if win32api.GetLastError() == winerror.ERROR_ALREADY_EXISTS:
        print("[!] HoneyFile Monitor is currently running silently in the background.")
        sys.exit(0)
    return mutex

# ================== MAIN ==================
if __name__ == "__main__":
    if not is_admin():
        # Cleanly relaunch as admin
        print("[*] Requesting Administrator Privileges to enable advanced SMB file tracking...")
        params = f'"{os.path.abspath(__file__)}" ' + ' '.join(sys.argv[1:])
        try:
            ctypes.windll.shell32.ShellExecuteW(None, "runas", sys.executable, params, None, 1)
        except Exception as e:
            print(f"[!] User denied UAC or failed to elevate: {e}")
        sys.exit(0)
    if "--foreground" not in sys.argv and os.path.basename(sys.executable).lower() == "python.exe":
        pythonw_path = sys.executable.replace("python.exe", "pythonw.exe")
        if os.path.exists(pythonw_path):
            import subprocess
            subprocess.Popen([pythonw_path, os.path.abspath(__file__), "--foreground"])
            sys.exit(0)

    _instance_mutex = prevent_multiple_instances()
    add_to_startup()

    os.makedirs(HONEY_DIR, exist_ok=True)
    os.makedirs(os.path.dirname(LOG_FILE), exist_ok=True)

    # Initialise tracked files setup
    for fname, fpath in HONEY_FILES.items():
        if not os.path.exists(fpath):
            try:
                ext = os.path.splitext(fname)[1].lower()
                base_text = f"CONFIDENTIAL DATA FILE - {fname}\nDo not share this file with unauthorized personnel."
                
                if ext == '.xlsx':
                    try:
                        import openpyxl
                        wb = openpyxl.Workbook()
                        ws = wb.active
                        ws.title = "Confidential Data"
                        ws['A1'] = base_text
                        wb.save(fpath)
                    except Exception as e:
                        print(f"[!] openpyxl failed for {fpath}, writing text fallback: {e}")
                        with open(fpath, "w") as f: f.write(base_text)
                
                elif ext == '.pdf':
                    try:
                        from fpdf import FPDF
                        pdf = FPDF()
                        pdf.add_page()
                        pdf.set_font("Arial", size=12)
                        pdf.cell(200, 10, txt=base_text, ln=1, align='C')
                        pdf.output(fpath)
                    except Exception as e:
                        print(f"[!] fpdf failed for {fpath}, writing text fallback: {e}")
                        with open(fpath, "w") as f: f.write(base_text)
                        
                elif ext == '.pptx':
                    try:
                        from pptx import Presentation
                        prs = Presentation()
                        slide = prs.slides.add_slide(prs.slide_layouts[0])
                        slide.shapes.title.text = "Confidential Data"
                        slide.placeholders[1].text = base_text
                        prs.save(fpath)
                    except Exception as e:
                        print(f"[!] python-pptx failed for {fpath}, writing text fallback: {e}")
                        with open(fpath, "w") as f: f.write(base_text)
                
                elif ext == '.zip':
                    import zipfile
                    with zipfile.ZipFile(fpath, 'w') as zf:
                        zf.writestr('readme.txt', base_text)
                
                elif ext == '.docx':
                    # A dummy text acts like a corrupted docx, to make it valid we need python-docx
                    # For now just write it if python-docx isn't installed
                    try:
                        import docx
                        doc = docx.Document()
                        doc.add_heading('Confidential Data', 0)
                        doc.add_paragraph(base_text)
                        doc.save(fpath)
                    except ImportError:
                        # docx module not required by user right now but fallback to write
                        with open(fpath, "w") as f: f.write(base_text)
                        
                else:
                    with open(fpath, "w") as f:
                        f.write(base_text)
                
                # Make sure "Everyone" has Read access so it shows up in SMB networks
                import subprocess
                subprocess.run(f'icacls "{fpath}" /grant "Everyone:(R)"', shell=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                        
            except Exception as e:
                print(f"[!] Could not create dummy file {fpath}: {e}")

        # Map initial hashes
        h = get_file_hash(fpath)
        ORIGINAL_HASHES[fname] = h
        last_hashes[fname] = h

    observer      = Observer()
    event_handler = HoneyFileHandler()
    observer.schedule(event_handler, path=HONEY_DIR, recursive=True)
    observer.start()

    print("=" * 50)
    print("  Project Cerberus Started")
    print("=" * 50)
    print(f"  Watching : {', '.join(HONEY_FILES.keys())}")
    print(f"  Log file : {LOG_FILE}")
    print("=" * 50)

    threading.Thread(target=monitor_drives,       args=(observer, event_handler), daemon=True).start()
    threading.Thread(target=process_copy_events,  daemon=True).start()
    threading.Thread(target=monitor_file_open,    daemon=True).start()
    threading.Thread(target=monitor_realtime_typing, daemon=True).start()
    threading.Thread(target=monitor_smb_sessions, daemon=True).start()
    threading.Thread(target=monitor_security_logs, daemon=True).start()

    try:
        while True: time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()